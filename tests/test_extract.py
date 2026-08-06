"""
Test dell'estrattore del corpus business.

Il caso che conta è il deck commerciale: le slide segnalate per ispezione visiva
sono quelle dove la promessa architetturale è disegnata invece che scritta, cioè
esattamente quelle che non devi saltare. Se lo script ti manda a guardare una
cartella che non ha prodotto, le salti.

  pytest tests/ -v
"""

from __future__ import annotations

import importlib.util
import locale
import subprocess
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
EXTRACT = ROOT / "skills" / "framework-capture" / "scripts" / "extract.py"


def _extract_module():
    """Importa lo script come modulo: `run()` va provata da sola, perché il difetto
    che conta sta nella decodifica dell'output del processo figlio."""
    spec = importlib.util.spec_from_file_location("extract_under_test", EXTRACT)
    mod = importlib.util.module_from_spec(spec)
    # In sys.modules prima di eseguirlo: con `from __future__ import annotations`
    # @dataclass risolve i tipi cercando il proprio modulo per nome, e non trovandolo
    # fallisce con un AttributeError che non ha niente a che vedere con l'estrazione.
    sys.modules[spec.name] = mod
    spec.loader.exec_module(mod)
    return mod


# Byte 0x9d: valido come continuazione UTF-8 di U+201D, non assegnato in cp1252.
# È il byte esatto su cui l'estrazione del corpus si è rotta in silenzio.
_CAMPIONE = "un cliente “attivo”".encode("utf-8")


def _il_locale_riproduce_il_difetto() -> bool:
    try:
        _CAMPIONE.decode(locale.getpreferredencoding(False))
    except (UnicodeDecodeError, LookupError):
        return True
    return False


@pytest.mark.skipif(not _il_locale_riproduce_il_difetto(),
                    reason="questo locale decodifica già UTF-8: il difetto non è "
                           "riproducibile qui, e un assert che passa sempre non protegge")
def test_child_output_is_read_as_utf8_and_not_as_the_system_codepage():
    """`pdftotext` emette UTF-8; il locale di Windows è cp1252. Decodificare
    l'output del figlio col locale solleva UnicodeDecodeError dentro il thread
    lettore di subprocess: l'eccezione non propaga, la pagina torna vuota, e
    finisce fra quelle «povere di testo».

    Lo strumento allora ti manda a guardare l'immagine di una pagina che di
    immagini non ne ha — cioè afferma una cosa falsa invece di tacere. Un
    estrattore che perde tre quarti di un documento senza dirlo è peggio di uno
    che si ferma, perché il corpus incompleto passa per completo."""
    extract = _extract_module()
    rc, out = extract.run([
        sys.executable, "-c",
        "import sys; sys.stdout.buffer.write(%r)" % (_CAMPIONE,),
    ])

    assert rc == 0, f"il figlio è fallito: {out!r}"
    assert "attivo" in out, (
        f"testo perso nella decodifica: {out!r}. La pagina risulterà vuota e "
        "verrà segnalata come contenuto grafico.")


pptx = pytest.importorskip("pptx", reason="serve python-pptx: pip install markitdown[pptx]")


@pytest.fixture
def deck(tmp_path: Path) -> Path:
    """Un deck con una slide povera di testo: l'architettura è un disegno."""
    from pptx import Presentation

    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text = "Piattaforma unificata"
    s.shapes.placeholders[1].text_frame.text = "Tre moduli, un'unica esperienza per il cliente"
    p.slides.add_slide(p.slide_layouts[5]).shapes.title.text = "Architettura"

    corpus = tmp_path / "corpus"
    corpus.mkdir()
    p.save(str(corpus / "deck.pptx"))
    return corpus


def _run(corpus: Path, out: Path) -> str:
    proc = subprocess.run(
        [sys.executable, str(EXTRACT), str(corpus), "-o", str(out)],
        capture_output=True, text=True, encoding="utf-8", errors="replace",
        env={**__import__("os").environ, "PYTHONIOENCODING": "utf-8"},
    )
    assert proc.returncode == 0, proc.stderr
    return proc.stdout


def test_the_folder_readme_is_not_ingested_as_a_business_document(deck: Path, tmp_path: Path):
    """Ogni cartella `corpus/` contiene un README che spiega cosa mettercì dentro.
    È scritto dal framework, non dal business: ingestarlo significa mettere in
    `ING.md` le nostre istruzioni con la forma di un'affermazione del cliente."""
    (deck / "README.md").write_text(
        "# Corpus\n\nI documenti del business che parlano solo di questo prodotto.\n",
        encoding="utf-8")

    out = tmp_path / "ingest-out"
    stdout = _run(deck, out)

    text = (out / "extract.md").read_text(encoding="utf-8")
    assert "README" not in text, "il README della cartella finisce fra il materiale da classificare"
    assert "README" not in stdout, f"il README viene contato come documento:\n{stdout}"


def test_output_survives_a_legacy_windows_console(deck: Path, tmp_path: Path):
    """Stessa console del validatore, stesso problema: uno strumento che va in
    crash quando ha qualcosa da dire è peggio di uno assente."""
    proc = subprocess.run(
        [sys.executable, str(EXTRACT), str(deck), "-o", str(tmp_path / "out")],
        capture_output=True, text=True, encoding="cp1252", errors="replace",
        env={**__import__("os").environ, "PYTHONIOENCODING": "cp1252"},
    )
    assert "UnicodeEncodeError" not in proc.stderr, \
        f"l'estrattore va in crash su una console non-UTF8:\n{proc.stderr[-600:]}"
    assert proc.returncode == 0, proc.stderr[-600:]


def test_pptx_is_extracted(deck: Path, tmp_path: Path):
    """Senza markitdown il deck produce zero blocchi e una nota morbida."""
    out = tmp_path / "ingest-out"
    _run(deck, out)
    text = (out / "extract.md").read_text(encoding="utf-8")
    assert "un'unica esperienza" in text, "il testo del deck non è stato estratto"
    assert "slide 1" in text, "manca la provenienza: senza, l'estrazione è inutilizzabile"


def test_does_not_send_you_to_a_directory_it_did_not_create(deck: Path, tmp_path: Path):
    """Un rimando a render/ vuoto insegna a ignorare la segnalazione."""
    out = tmp_path / "ingest-out"
    stdout = _run(deck, out)

    rendered = list((out / "render").glob("*")) if (out / "render").exists() else []
    if not rendered:
        assert "render/" not in stdout, (
            "lo script rimanda a render/ ma non ha prodotto nessuna immagine:\n" + stdout)


def test_unrenderable_pages_name_the_file_to_open(deck: Path, tmp_path: Path):
    """Se non può rasterizzare, deve dire quale file aprire e a quale slide."""
    out = tmp_path / "ingest-out"
    stdout = _run(deck, out)

    rendered = list((out / "render").glob("*")) if (out / "render").exists() else []
    if not rendered:
        assert "deck.pptx" in stdout and "2" in stdout, (
            "non dice quale file aprire e dove:\n" + stdout)
